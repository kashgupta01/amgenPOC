"""
Document extraction pipeline.

Supported formats: PDF (.pdf), Word (.docx), PowerPoint (.pptx), Excel (.xlsx/.xls/.csv)

Usage:
    blocks = process_file("path/to/doc.pdf")
"""

import pathlib
from dataclasses import dataclass


#Data model ────────────────────────────────────────────────────────────────
#The extracted text is stored in a simple dataclass called ExtractedBlock, which includes the text content, the source file name, the type of source (pdf, docx, pptx, excel), and the location within the file (e.g., page number, slide number, sheet name and row).
#Why do we need a dataclass? It provides a structured way to represent the extracted information, making it easier to manage and use in downstream processing. Each ExtractedBlock contains not just the text, but also metadata about where it came from, which can be crucial for context when analyzing the data later on.
#How do dataclasses work? The @dataclass decorator automatically generates special methods like __init__, __repr__, and __eq__ based on the class attributes we define. This means we can easily create instances of ExtractedBlock without having to write boilerplate code for initialization and representation, making our code cleaner and more maintainable.
@dataclass
class ExtractedBlock:
    text: str
    source_file: str
    source_type: str        # "pdf" | "docx" | "pptx" | "excel"
    location: str           # e.g. "page 3", "slide 2", "sheet Sales row 4-10"


# ── Extractors ────────────────────────────────────────────────────────────────

#Each supported file type has a corresponding extractor function that handles the specific logic for reading and extracting text from that format. These functions use popular libraries like fitz for PDFs, python-docx for Word documents, python-pptx for PowerPoint presentations, and pandas for Excel files. The extractors return a list of tuples containing the extracted text and its location within the file, which are then used to create ExtractedBlock instances in the main processing function.
#what is def _extract_pdf(path: pathlib.Path) -> list[tuple[str, str]]? This is a function definition for _extract_pdf, which takes a file path as input and returns a list of tuples. Each tuple contains a string of extracted text and a string indicating the location of that text within the PDF (e.g., "page 1"). The function uses the fitz library to open the PDF document, iterate through its pages, extract text from each page, and store it along with the page number in the results list.
#why do we need to specify input and output? Specifying input and output types in the function definition (using type hints) helps improve code readability and maintainability. It provides clear information about what kind of data the function expects and what it will return, which can help catch errors early and make it easier for other developers (or your future self) to understand how to use the function correctly.
def _extract_pdf(path: pathlib.Path) -> list[tuple[str, str]]:
    import fitz
    import re

    # Patterns that mark noise: figure/table captions, footnote markers, reference lists
    _NOISE = re.compile(
        r"^(fig(ure)?s?\.?\s*\d|table\s*\d|supplementary|©|doi:|https?://|\[\d+\]|\d+\s+[A-Z][\w\s]{0,30}et al)",
        re.IGNORECASE,
    )

    results = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, 1):
            page_h = page.rect.height
            footnote_threshold = page_h * 0.88  # bottom 12% = footnote zone

            kept = []
            #what is page.get_text("blocks")? This method retrieves the text from the PDF page in the form of blocks, where each block is a tuple containing the coordinates of the block (x0, y0, x1, y1), the text content of the block, the block number, and the block type. The block type indicates whether the block is text or an image, allowing us to filter out non-text elements during extraction.
            for block in page.get_text("blocks"):
                # block: (x0, y0, x1, y1, text, block_no, block_type)
                # block_type 1 = image — skip entirely
                #why block_type 1 is an image? In the fitz library, when you call page.get_text("blocks"), it returns a list of blocks where each block is a tuple containing information about the block's position, text content, block number, and block type. The block type is an integer that indicates the type of content in the block. A block type of 1 specifically indicates that the block contains an image rather than text. This allows us to easily filter out image blocks during the extraction process, ensuring that we only keep textual content for our knowledge graph construction.
                #what are the block types? The block types in the fitz library are typically defined as follows:
                #0: Text block - contains text content that can be extracted.
                #1: Image block - contains an image, which is not text and should be skipped for our purposes.
                #2: Vector block - contains vector graphics, which may include lines, shapes, or other non-text elements that we also want to skip. 
                
                if block[6] == 1 or block[6] == 2:
                    continue
                # skip blocks that start in the footnote zone
                y0 = block[1]
                if y0 > footnote_threshold:
                    continue
                text = block[4].strip()
                if not text:
                    continue
                # skip figure/table captions and reference noise
                #what is this line doing? This line uses a regular expression to check if the first line of the extracted text block matches common patterns that indicate noise, such as figure or table captions, supplementary material references, copyright symbols, DOIs, URLs, or reference list entries. If the first line matches any of these patterns, the block is considered noise and is skipped from being included in the final results. This helps to reduce irrelevant information in the extracted data, making it more focused and useful for building the knowledge graph.
                #where is it iterating through each line of the text block? The line is splitting the extracted text block into individual lines using the splitlines() method, and then it checks only the first line (first_line) against the regular expression pattern to determine if it should be considered noise. If the first line matches the noise pattern, the entire block is skipped. This approach assumes that if the first line of a block is noise, the rest of the block is likely to be noise as well, which helps to efficiently filter out irrelevant content from the PDF.
                first_line = text.splitlines()[0]
                if _NOISE.match(first_line):
                    continue
                kept.append(text)

            combined = "\n".join(kept).strip()
            if combined:
                results.append((combined, f"page {i}"))
    return results


def _extract_docx(path: pathlib.Path) -> list[tuple[str, str]]:
    from docx import Document
    doc = Document(path)
    results = []
    #iterate through each paragraph and extract text along with its location 
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if text:
            results.append((text, f"paragraph {i}"))
    return results

#why is there an underscore before the function names? The underscore before the function names (e.g., _extract_pdf) is a common convention in Python to indicate that these functions are intended for internal use within the module and are not part of the public API. It signals to other developers that these functions are meant to be private and should not be accessed directly from outside the module, helping to encapsulate the implementation details and maintain a cleaner interface.

def _extract_pptx(path: pathlib.Path) -> list[tuple[str, str]]:
    #
    from pptx import Presentation
    prs = Presentation(path)
    results = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            # Extract text from all shapes in the slide that have text frames and are not empty
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        # Combine all extracted text from the slide into a single string and store it with the slide number as location
        combined = "\n".join(texts)
        # Only add to results if there is any text extracted from the slide
        if combined:
            results.append((combined, f"slide {i}"))
    return results

#Excel files can have multiple sheets, and each sheet can have multiple rows of data. The _extract_excel function uses pandas to read the Excel file, handling both .xlsx and .csv formats. It iterates through each sheet (or the single sheet in case of CSV), and for each row that contains non-empty values, it combines the values into a single string and stores it along with the sheet name and row number as the location. This allows us to capture structured data from Excel files in a way that can be easily referenced later on.
def _extract_excel(path: pathlib.Path) -> list[tuple[str, str]]:
    import pandas as pd
    suffix = path.suffix.lower()
    frames = {"sheet 1": pd.read_csv(path)} if suffix == ".csv" else pd.read_excel(path, sheet_name=None)
    results = []
    for sheet_name, df in frames.items():
        for i, row in df.dropna(how="all").iterrows():
            text = " | ".join(str(v) for v in row.values if str(v).strip())
            if text:
                results.append((text, f"sheet '{sheet_name}' row {i + 1}"))
    return results

#Mapping of file extensions to their corresponding extractor functions. This allows the main processing function to dynamically select the appropriate extractor based on the file type, making the code more modular and easier to maintain. If a new file type needs to be supported in the future, we can simply add a new extractor function and update this mapping without having to change the core logic of the process_file function.
_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_excel,
    ".xls":  _extract_excel,
    ".csv":  _extract_excel,
}


#Entry point ───────────────────────────────────────────────────────────────

#The process_file function serves as the main entry point for extracting text from any supported file type. It takes a file path as input, determines the file type based on its extension, and then calls the appropriate extractor function from the _EXTRACTORS mapping. The extracted text and its location are then used to create a list of ExtractedBlock instances, which can be easily used in downstream processing or analysis. This function abstracts away the details of how each file type is processed, providing a simple interface for extracting text regardless of the source format.
def process_file(path: str | pathlib.Path) -> list[ExtractedBlock]:
    """Extract text from any supported file and return a flat list of ExtractedBlocks."""
    path = pathlib.Path(path)
    suffix = path.suffix.lower()

    #determines the extension of the file and looks up the corresponding extractor function in the _EXTRACTORS dictionary. If the file type is not supported (i.e., there is no corresponding extractor), it raises a ValueError with a message indicating the unsupported file type and listing the supported types. This ensures that the function fails gracefully when given an unsupported file, providing clear feedback to the user about what went wrong and how to fix it.
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix!r}. Supported: {list(_EXTRACTORS)}")
    
    #calls the appropriate extractor function to get a list of tuples containing the extracted text and its location. It then creates and returns a list of ExtractedBlock instances, where each block contains the extracted text, the source file name, the source type (derived from the file extension), and the location within the file. This structured representation allows for easy handling of the extracted data in later stages of processing or analysis.
    return [
        ExtractedBlock(text=text, source_file=path.name, source_type=suffix.lstrip("."), location=location)
        for text, location in extractor(path)
    ]


