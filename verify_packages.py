"""
Package verification and comparison.

Sections
--------
1. Existing smoke-tests      ->  verification_results.txt
2. PDF package comparison    ->  comparison_pdf.txt
3. Word-docx comparison      ->  comparison_docx.txt
4. PPTX comparison           ->  comparison_pptx.txt

Required packages (install if missing):
    pip install pdfplumber "pdfminer.six" pypdf tika mammoth docx2txt

Note: Apache Tika requires Java. It downloads tika-server.jar on first run (~60 s).
"""

import pathlib
import sys

SAMPLE = pathlib.Path(__file__).parent / "src" / "data_plane" / "sample_data"
DATA   = pathlib.Path(__file__).parent / "src" / "data_plane" / "data"
OUTPUT = pathlib.Path(__file__).parent / "verification_results.txt"

PREVIEW = 2000  # characters shown per library in comparison files


# ── Tee helper ────────────────────────────────────────────────────────────────

class _Tee:
    """Writes to both stdout and a file simultaneously."""
    def __init__(self, file):
        self._file = file
        self._stdout = sys.stdout

    def write(self, data):
        self._stdout.write(data)
        self._file.write(data)

    def flush(self):
        self._stdout.flush()
        self._file.flush()


# ── Existing smoke-tests ──────────────────────────────────────────────────────

_out_file = open(OUTPUT, "w", encoding="utf-8")
_tee = _Tee(_out_file)
sys.stdout = _tee

#PyMuPDF
print("=" * 60)
print("1. PyMuPDF — Bioinformatics PDF")
print("=" * 60)

import fitz

pdf_path = SAMPLE / "Bioinformatics analysis of the role of RNA modification regulators in polycystic ovary syndrome.pdf"

with fitz.open(pdf_path) as doc:
    meta = doc.metadata
    print(f"  Title  : {meta.get('title') or '(not embedded)'}")
    print(f"  Author : {meta.get('author') or '(not embedded)'}")
    first_page_text = doc[0].get_text()[:500].strip()
    print(f"  Page 1 preview:\n{first_page_text}\n")


#pandas
print("=" * 60)
print("2. pandas — iris_data.csv")
print("=" * 60)

import pandas as pd

df = pd.read_csv(SAMPLE / "iris_data.csv")
print(f"  Shape   : {df.shape[0]} rows × {df.shape[1]} columns")
print(f"  Columns : {list(df.columns)}")
print(f"  Head:\n{df.head(3).to_string(index=False)}\n")


#python-pptx
print("=" * 60)
print("3. python-pptx — DS340 Discussion1.pptx")
print("=" * 60)

from pptx import Presentation

pptx_path = SAMPLE / "DS340 – Discussion1.pptx"
prs = Presentation(pptx_path)
print(f"  Slides : {len(prs.slides)}")

for i, slide in enumerate(list(prs.slides)[:3], 1):
    texts = [
        shape.text.strip()
        for shape in slide.shapes
        if shape.has_text_frame and shape.text.strip()
    ]
    print(f"  Slide {i}: {' | '.join(texts[:3]) or '(no text)'}")
print()


#python-docx
print("=" * 60)
print("4. python-docx — Flaws_In_Clinical_Methodology_Kashish.docx")
print("=" * 60)

from docx import Document

doc_path = SAMPLE / "Flaws_In_Clinical_Methodology_Kashish.docx"
doc = Document(doc_path)

core = doc.core_properties
print(f"  Title  : {core.title or '(not set)'}")
print(f"  Author : {core.author or '(not set)'}")

non_empty = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
print(f"  Paragraphs (non-empty): {len(non_empty)}")
print(f"  First paragraph: {non_empty[0][:200] if non_empty else '(empty)'}\n")

print("=" * 60)
print("All packages verified successfully.")
print("=" * 60)

sys.stdout = _tee._stdout
_out_file.close()
print(f"\nResults saved to: {OUTPUT}")


# ── Comparison helpers ────────────────────────────────────────────────────────

SEP  = "=" * 70
DASH = "─" * 70


def _write_comparison(title, sample_path, extractors, out_path):
    lines = [
        SEP,
        f"{title} — PACKAGE COMPARISON",
        f"Sample file : {sample_path.name}",
        SEP, "",
    ]
    for lib_name, fn in extractors:
        lines += [DASH, f"LIBRARY: {lib_name}", DASH]
        try:
            text = fn(sample_path).strip()
            snippet = text[:PREVIEW] + (" ...[truncated]" if len(text) > PREVIEW else "")
            lines.append(snippet)
        except Exception as exc:
            lines.append(f"[ERROR — {exc}]")
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"  {title:<6}  ->  {out_path.name}")


# ── PDF extractors ────────────────────────────────────────────────────────────

def _pdf_pymupdf(path):
    import fitz
    return "\n".join(page.get_text() for page in fitz.open(path))

def _pdf_pdfplumber(path):
    import pdfplumber
    with pdfplumber.open(path) as pdf:
        return "\n".join(p.extract_text() or "" for p in pdf.pages)

def _pdf_pdfminer(path):
    from pdfminer.high_level import extract_text
    return extract_text(str(path))

def _pdf_pypdf(path):
    from pypdf import PdfReader
    return "\n".join(p.extract_text() for p in PdfReader(str(path)).pages)

def _pdf_tika(path):
    from tika import parser
    return parser.from_file(str(path)).get("content") or ""


# ── DOCX extractors ───────────────────────────────────────────────────────────

def _docx_python_docx(path):
    from docx import Document
    return "\n".join(p.text for p in Document(path).paragraphs if p.text.strip())

def _docx_mammoth(path):
    import mammoth
    with open(path, "rb") as f:
        return mammoth.extract_raw_text(f).value

def _docx_docx2txt(path):
    import docx2txt
    return docx2txt.process(str(path))


# ── PPTX extractors ───────────────────────────────────────────────────────────

def _pptx_python_pptx(path):
    from pptx import Presentation
    slides = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        texts = [s.text.strip() for s in slide.shapes if s.has_text_frame and s.text.strip()]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)

def _pptx_tika(path):
    from tika import parser
    return parser.from_file(str(path)).get("content") or ""


# ── Run comparisons ───────────────────────────────────────────────────────────

print("\nRunning package comparisons ...")
_write_comparison(
    "PDF",
    SAMPLE / "Bioinformatics analysis of the role of RNA modification regulators in polycystic ovary syndrome.pdf",
    [
        ("PyMuPDF (fitz)", _pdf_pymupdf),
        ("pdfplumber",     _pdf_pdfplumber),
        ("pdfminer.six",   _pdf_pdfminer),
        ("pypdf",          _pdf_pypdf),
        ("Apache Tika",    _pdf_tika),
    ],
    pathlib.Path(__file__).parent / "comparison_pdf.txt",
)

_write_comparison(
    "DOCX",
    SAMPLE / "Flaws_In_Clinical_Methodology_Kashish.docx",
    [
        ("python-docx", _docx_python_docx),
        ("mammoth",     _docx_mammoth),
        ("docx2txt",    _docx_docx2txt),
    ],
    pathlib.Path(__file__).parent / "comparison_docx.txt",
)

_write_comparison(
    "PPTX",
    SAMPLE / "DS340 – Discussion1.pptx",
    [
        ("python-pptx", _pptx_python_pptx),
        ("Apache Tika", _pptx_tika),
    ],
    pathlib.Path(__file__).parent / "comparison_pptx.txt",
)

print("\nDone.")
